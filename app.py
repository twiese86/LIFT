import os
import json
from io import BytesIO
from flask import Flask, request, jsonify, Response, session, send_file
from flask_cors import CORS
import google.generativeai as genai
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# -----------------------------
# LIFT: Web App with Chat UI + Conversational Memory
# -----------------------------

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}})

app.secret_key = os.getenv("FLASK_SECRET_KEY", "change-me-in-production")

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
if not GEMINI_API_KEY:
    raise RuntimeError("Set the GEMINI_API_KEY environment variable before running the app.")

genai.configure(api_key=GEMINI_API_KEY)

MODEL_NAME = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")
model = genai.GenerativeModel(MODEL_NAME)

app.config["MAX_CONTENT_LENGTH"] = 10 * 1024 * 1024  # 10 MB uploads

MAX_HISTORY_TURNS = 6
USER_SNIPPET_CHARS = 600
ASSISTANT_SNIPPET_CHARS = 2000

# ── PII scrubbing (lightweight regex — no heavy dependencies) ────────────────
import re

_PII_PATTERNS = [
    (re.compile(r'\b[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}\b'), '[EMAIL]'),
    (re.compile(r'\b(?:\+?1[\s\-.]?)?\(?\d{3}\)?[\s\-.]?\d{3}[\s\-.]?\d{4}\b'), '[PHONE]'),
    (re.compile(r'\b\d{3}-\d{2}-\d{4}\b'), '[SSN]'),
    (re.compile(r'\b(?:\d[ \-]?){13,16}\b'), '[CARD]'),
    (re.compile(r'\b[A-Z][a-z]+ [A-Z][a-z]+\b'), '[NAME]'),   # simple two-word proper name
]

def scrub_pii(text: str) -> str:
    """Redact common PII patterns before sending to the LLM."""
    for pattern, replacement in _PII_PATTERNS:
        text = pattern.sub(replacement, text)
    return text

# Mapping of use cases to system context
USE_CASES = {

    "none": (
        "General instructional support. Help faculty with any course design, content development, "
        "or pedagogical question. Apply sound instructional design principles and cite relevant "
        "learning theory where useful."
    ),

    # ── Existing use cases (expanded with learning theory) ────

    "uc1": (
        "Context: Rapid Course Material Development for New Faculty.\n\n"
        "Apply Gagné's Nine Events of Instruction when structuring every output: "
        "open with an attention-gaining element (provocative question, compelling statistic, current-event hook); "
        "state clear, measurable objectives using action verbs aligned to verbal information, intellectual skills, "
        "cognitive strategies, attitudes, or motor skills; stimulate recall of prior learning; "
        "present content with worked examples and encoding guidance (mnemonics, analogies, concept maps); "
        "elicit practice; provide feedback with model answers and rubrics; assess with both formative and "
        "summative items; and include transfer activities across varied contexts.\n\n"
        "Apply Schema Theory: organize content hierarchically (foundational → complex), explicitly connect "
        "new concepts to prior knowledge, address common misconceptions with refutational explanations, and "
        "progress through acclimation → competence → proficiency.\n\n"
        "Apply Information Processing Theory: chunk content into 5–7 item units, minimize cognitive load "
        "through clear scaffolding, provide multiple encoding strategies (verbal, visual, kinesthetic), "
        "and design spaced retrieval practice.\n\n"
        "Output format: weekly module structure; each lesson contains objectives, prerequisite review, "
        "content presentation, guided practice, independent practice, assessment, and application activities. "
        "State which Gagné outcome type each objective addresses and include an authentic unifying context."
    ),

    "uc2": (
        "Context: Creating Accessible & Multi-Format Learning Materials.\n\n"
        "Apply Information Processing principles across modalities: for visual learners provide diagrams "
        "and concept maps with integrated text (avoid split-attention); for auditory learners generate "
        "narrated content that explains rather than reads verbatim; apply Dual Coding by pairing verbal "
        "and visual representations.\n\n"
        "Reduce Cognitive Load (CLT): avoid the split-attention and redundancy effects; follow the "
        "modality and segmenting principles; use worked examples for complex procedures.\n\n"
        "Implement Universal Design for Learning (UDL): multiple means of representation (text, audio, "
        "visual); multiple means of action/expression; multiple means of engagement with varied contexts "
        "connecting to diverse student backgrounds. If the tool cannot produce multimedia directly, "
        "suggest complementary AI tools (e.g., NotebookLM for audio).\n\n"
        "Accessibility standards: WCAG 2.1 AA compliance; alt-text describing function and meaning; "
        "transcripts with timestamps; simplified text at 8th-grade reading level alongside standard versions; "
        "screen-reader-optimized HTML with proper heading hierarchy; high-contrast typography.\n\n"
        "Output: master content with semantic structure, auto-generated alt-text, transcripts, "
        "simplified and standard versions, and self-assessment reflection prompts in every format."
    ),

    "uc3": (
        "Context: Automated Assessment Creation & Rubric Development.\n\n"
        "Align every assessment item to Gagné's five learning outcome types: verbal information (recall/recognition), "
        "intellectual skills (application, problem-solving), cognitive strategies (metacognitive prompts), "
        "attitudes (scenario-based value judgments), motor skills (procedural checklists).\n\n"
        "Apply Bloom's Taxonomy hierarchically across all six levels—Remember, Understand, Apply, Analyze, "
        "Evaluate, Create—with explicit emphasis on the upper three levels (analyze, evaluate, create). "
        "Distribute items across levels intentionally and label each item with its Bloom's level.\n\n"
        "Incorporate Schema Theory: include items that reveal misconceptions, differentiate novice from "
        "expert thinking, and require far transfer (not just near transfer) to new contexts.\n\n"
        "Support Constructivist and Authentic Assessment: design tasks that mirror real-world applications; "
        "include performance-based, portfolio, and presentation options; provide structured peer-review "
        "protocols and self-assessment rubrics with metacognitive reflection prompts.\n\n"
        "Rubric design requirements: clearly defined criteria aligned to objectives; 3–5 performance levels "
        "with descriptive (not merely evaluative) language; observable behavioral indicators at each level; "
        "developmental progression from novice to expert; feed-forward guidance showing how to improve.\n\n"
        "Output: question banks organized by outcome type, Bloom's level, difficulty, and topic; "
        "multiple assessment versions; answer keys with explanations and common error patterns; "
        "analytic and holistic rubrics with anchor examples."
    ),

    "uc4": (
        "Context: Flipped Classroom Content Generation.\n\n"
        "Design pre-class materials using Information Processing & Schema Theory: segment content into "
        "5–10 minute chunks; provide advance organizers showing how new material connects to prior knowledge; "
        "apply Dual Coding (visuals + verbal explanations); include worked examples with think-aloud "
        "commentary; embed formative questions requiring active processing; focus on Bloom's "
        "Remember/Understand levels.\n\n"
        "Design in-class materials using Constructivism & Situated Learning: target Bloom's Apply, Analyze, "
        "Evaluate, and Create levels; use authentic real-world problems; include collaborative structures "
        "(think-pair-share, jigsaw, peer instruction); scaffold problem-solving by gradually fading support; "
        "shift the instructor role from information delivery to facilitation and coaching.\n\n"
        "Apply the 4C/ID Cognitive Load model: low cognitive load in pre-class; complex integrated tasks "
        "with instructor support in-class; whole-task practice; progressive complexity; just-in-time "
        "procedural information delivered exactly when needed.\n\n"
        "Support Self-Regulated Learning: generate self-check quizzes with immediate feedback, learning "
        "journal prompts, goal-setting templates, and metacognitive reflection questions.\n\n"
        "Output — Pre-Class: video scripts (8–12 min per concept), reading guides with embedded questions, "
        "advance organizer, self-assessment quiz (5–10 items with explanatory feedback), note-taking "
        "template, 'muddy point' reflection form.\n"
        "Output — In-Class: graduated problem sets, group activity protocols with defined roles, case studies "
        "at multiple Bloom's levels, peer-teaching activities, formative checkpoint items, instructor "
        "facilitation guide with common misconceptions.\n"
        "Output — Synthesis: worksheets connecting pre-class concepts to in-class applications, cumulative "
        "integration projects, self-assessment rubrics for skill development."
    ),

    "uc5": (
        "Context: Cross-Disciplinary Course Revision & Modernization.\n\n"
        "Begin with a learning-theory-based content analysis: classify existing objectives by Gagné outcome "
        "type and Bloom's level; identify cognitive overload points; map where content assumes prerequisite "
        "understanding vs. builds on it; evaluate transfer potential.\n\n"
        "Update content using Information Processing: add contemporary, culturally relevant examples; "
        "insert graphic organizers and mnemonics; design spaced review and interleaved practice; direct "
        "attention through relevance statements.\n\n"
        "Modernize using Constructivist principles: replace decontextualized exercises with authentic "
        "real-world problems; add inquiry-based elements; incorporate diverse voices and historically "
        "marginalized perspectives; provide student choice in demonstrating understanding.\n\n"
        "Apply Situated Learning: design activities mirroring how practitioners in the discipline actually "
        "think and work; create scaffolded entry into disciplinary discourse (legitimate peripheral "
        "participation); include collaborative knowledge-building to replace isolated individual work.\n\n"
        "Integrate Digital Literacy (Connectivism): add source-credibility evaluation activities, "
        "primary source analysis, and tasks requiring synthesis across diverse information networks.\n\n"
        "Address Equity & Inclusion: update examples to reflect diverse student backgrounds; provide "
        "multiple entry points; include counter-narratives; frame diverse backgrounds as assets.\n\n"
        "Output: gap analysis report mapping current content to learning theories with prioritized "
        "recommendations; revised syllabus with modernized measurable objectives; updated reading list "
        "including OERs; modernized assessments with formative/summative balance; implementation guide "
        "with learning-theory rationale, scaffolding plan, and technology integration recommendations."
    ),

    # ── New use cases ─────────────────────────────────────────

    "uc6": (
        "Context: Converting Notes or Readings into Learning Theory-Based Lecture Presentations.\n\n"
        "Structure every presentation using Gagné's Nine Events: "
        "Slides 1–2 gain attention (provocative question, surprising statistic, current-event hook) and "
        "state measurable objectives with action verbs; Slide 3 stimulates recall via a concept map or "
        "formative question connecting to prerequisites; Content slides present stimulus and provide "
        "learning guidance.\n\n"
        "Apply Cognitive Load Theory / Multimedia Learning Principles: "
        "Coherence — remove extraneous decorative elements; "
        "Signaling — use arrows, color, bold to highlight key information; "
        "Redundancy — use brief bullet points on slides, not full narration text; "
        "Spatial Contiguity — place text labels near corresponding graphics; "
        "Segmenting — break complex content across multiple slides with learner control; "
        "Modality — pair visual slides with narration script, not identical on-screen text.\n\n"
        "Support Information Processing: max 5–7 items per slide; progressive disclosure; dual coding "
        "(verbal + visual); 2–3 concrete examples per abstract concept; visual hierarchy showing "
        "relationships; ample white space.\n\n"
        "Build schemas: include an advance organizer slide; progress from general to specific; add "
        "misconception-refutation slides; close sections with conceptual integration slides.\n\n"
        "Embed active learning every 10–15 minutes: think-pair-share prompts, polling questions, "
        "problem-solving scenarios, reflection prompts, and application mini-cases. Include presenter "
        "notes indicating pause points for each activity.\n\n"
        "Final slides: visual summary of key takeaways; application scenarios; preview of upcoming "
        "topics; 3–5 self-assessment questions; resources slide.\n\n"
        "For each slide provide: (1) minimal slide content — brief bullets/visuals; "
        "(2) presenter notes — full narration script, learning-theory rationale, timing, common student "
        "questions, transition statements; (3) accessibility notes — alt-text, pronunciation guides, "
        "contrast information.\n\n"
        "Supplementary materials: handout version with note-taking space; accessible full-text alternative; "
        "study guide with key terms and self-test questions; concept map of all relationships covered."
    ),

    "uc7": (
        "Context: Faculty AI Literacy & Professional Development.\n\n"
        "Support a faculty cohort in building structured AI competency aligned with the IITG objective "
        "to 'develop instructional talent.' Generate guided AI prompt templates that faculty can "
        "practice and adapt; embed Bloom's Taxonomy alignment tools so participants can evaluate "
        "AI-generated materials against pedagogical standards; design reflection + revision workflows "
        "that ask faculty to critique, improve, and document AI outputs; create peer-sharing "
        "frameworks (discussion protocols, annotation guides) that build a sustainable community of practice.\n\n"
        "Ground all outputs in adult learning principles (andragogy): connect new AI skills to faculty "
        "members' existing course contexts; provide immediate applicability (faculty work on a real "
        "module they will teach); build in self-directed choice (faculty choose their use case); "
        "include formative self-assessment so participants can track their own AI confidence growth "
        "(measurable pre/post).\n\n"
        "Output: structured workshop sequence with learning objectives at each stage; prompt-template "
        "library covering all LIFT use cases with explanatory annotations; Bloom's alignment "
        "checklist for reviewing AI-generated content; reflection journal prompts; peer-review "
        "protocol for sharing and evaluating AI-assisted course modules; community-of-practice "
        "facilitation guide for ongoing cross-campus collaboration."
    ),

    "uc8": (
        "Context: Cross-Campus Collaborative Course Development.\n\n"
        "Support faculty from multiple SUNY campuses (e.g., Empire State University and Cobleskill College) "
        "in co-developing shared or interdisciplinary course materials. Generate shared rubric frameworks "
        "that can be harmonized across institutional contexts; produce OER-based content that can be "
        "freely remixed by contributors at each campus; create structured version-comparison outputs "
        "that make it easy to identify differences between campus-specific adaptations; support "
        "interdisciplinary integration by mapping learning objectives from two or more disciplines "
        "onto a unified framework.\n\n"
        "Apply Situated Learning and Communities of Practice: design materials that invite legitimate "
        "peripheral participation (new faculty or campus partners can contribute incrementally); "
        "build in negotiation and argumentation structures (collaborative rubric review, joint "
        "objective-setting protocols); make disciplinary norms and practices explicit so contributors "
        "from different fields can work across boundaries.\n\n"
        "Output: shared course outline template with fields for campus-specific customization; "
        "harmonized rubric with common criteria and campus-specific performance descriptors; "
        "OER reading list with remixing notes; interdisciplinary concept-mapping activity; "
        "collaboration workflow guide covering roles, review cycles, and version-tracking conventions."
    ),

    "uc9": (
        "Context: AI-Assisted Pedagogical Research Design.\n\n"
        "Support faculty researchers in using LIFT as a research instrument for studying instructional "
        "effectiveness. Generate paired content versions (Version A / Version B) for A/B comparison "
        "studies — e.g., AI-generated vs. manually created materials, or two different instructional "
        "approaches — with clearly documented generation parameters for each version; produce structured "
        "rubric comparison frameworks that allow researchers to evaluate materials against the same "
        "criteria; auto-generate case-study documentation capturing design decisions, theoretical "
        "rationale, and implementation notes suitable for inclusion in a scholarly write-up.\n\n"
        "Align with SUNY's research dissemination pillar: outputs should be formatted to support "
        "conference presentations (e.g., SUNY CIT), journal submissions, or IRB-ready study protocols. "
        "Include a methodology section template describing how LIFT was used, which use case context "
        "was applied, and what prompt parameters were set.\n\n"
        "Output: A/B content pair with documented generation parameters; comparative rubric with "
        "inter-rater reliability guidance; case study template (background, method, materials, "
        "analysis plan, findings placeholder); presentation-ready summary slide outline; "
        "suggested engagement-analytics data-collection points tied to each content version."
    ),
}

@app.route("/", methods=["GET"])
def ui():
    html = f"""
    <!DOCTYPE html>
    <html lang="en">
    <head>
      <meta charset="UTF-8">
      <meta name="viewport" content="width=device-width, initial-scale=1.0">
      <title>LIFT Tool</title>
      <style>
        * {{ box-sizing: border-box; }}
        body, p, div, h1, h2, h3, h4, h5, h6 {{ margin: 0; padding: 0; }}
        body {{ font-family: system-ui, -apple-system, Segoe UI, Roboto, Arial, sans-serif; background: #f3f4f6; color: #111827; }}
        .app {{ min-height: 100vh; display: flex; align-items: stretch; justify-content: center; padding: 1.5rem; }}
        .card {{ background: #ffffff; width: 100%; max-width: 960px; border-radius: 18px; box-shadow: 0 10px 30px rgba(15,23,42,0.08); display: flex; flex-direction: column; overflow: hidden; }}
        .header {{ padding: 1.25rem 1.5rem 0.75rem; border-bottom: 1px solid #e5e7eb; }}
        .title-row {{ display: flex; align-items: center; gap: 0.75rem; }}
        .pill {{ width: 32px; height: 32px; border-radius: 999px; background: linear-gradient(135deg, #6366f1, #06b6d4); display: flex; align-items: center; justify-content: center; color: white; font-weight: 700; font-size: 0.9rem; }}
        h1 {{ font-size: 1.1rem; }}
        .subtitle {{ margin-top: 0.4rem; font-size: 0.9rem; color: #6b7280; }}
        .body {{ display: flex; flex-direction: column; gap: 0.75rem; padding: 0.75rem 1.5rem 1rem; min-height: 0; flex: 1; }}
        .chat-window {{ border-radius: 12px; border: 1px solid #e5e7eb; background: #f3f4f6; padding: 0.75rem; overflow-y: auto; flex: 1; max-height: 60vh; }}
        .chat-window::-webkit-scrollbar {{ width: 6px; }}
        .chat-window::-webkit-scrollbar-thumb {{ background: #d1d5db; border-radius: 999px; }}
        .message {{ width: 100%; display: flex; align-items: flex-start; margin-bottom: 0.75rem; gap: 0.5rem; }}
        .message.assistant {{ flex-direction: row; }}
        .message.user {{ flex-direction: row-reverse; }}
        .avatar {{ width: 32px; height: 32px; border-radius: 999px; flex-shrink: 0; display: flex; align-items: center; justify-content: center; font-size: 0.75rem; font-weight: 600; }}
        .assistant .avatar {{ background: #eef2ff; color: #3730a3; }}
        .user .avatar {{ background: #ecfeff; color: #0f766e; }}
        .bubble {{ max-width: 60%; display: inline-block; margin: 0 6px; padding: 0.75rem 1rem; border-radius: 14px; font-size: 0.9rem; line-height: 1.45; white-space: normal; word-break: break-word; }}
        .assistant .bubble {{ background: #ffffff; border: 1px solid #e5e7eb; }}
        .user .bubble {{ background: #111827; color: #f9fafb; }}
        .name {{ font-size: 0.75rem; font-weight: 600; margin-bottom: 0.15rem; opacity: 0.8; }}
        .bubble-body {{ font-size: 0.9rem; text-align: left; }}
        .typing-dots {{ display: inline-flex; gap: 3px; align-items: center; }}
        .dot {{ width: 4px; height: 4px; border-radius: 999px; background: #9ca3af; animation: blink 1.4s infinite both; }}
        .dot:nth-child(2) {{ animation-delay: 0.2s; }}
        .dot:nth-child(3) {{ animation-delay: 0.4s; }}
        @keyframes blink {{ 0%, 80%, 100% {{ opacity: 0.3; }} 40% {{ opacity: 1; }} }}
        .form-wrapper {{ border-radius: 12px; border: 1px solid #e5e7eb; padding: 0.75rem 0.9rem 0.9rem; background: #ffffff; display: flex; flex-direction: column; gap: 0.5rem; }}
        .ppt-wrapper {{ border-radius: 12px; border: 1px solid #c7d2fe; padding: 0.75rem 0.9rem 0.9rem; background: #eef2ff; display: flex; flex-direction: column; gap: 0.5rem; }}
        .ppt-wrapper label {{ color: #3730a3; }}
        .section-title {{ font-size: 0.75rem; font-weight: 700; color: #6b7280; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 0.1rem; }}
        .ppt-section-title {{ font-size: 0.75rem; font-weight: 700; color: #3730a3; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 0.1rem; }}
        .form-row {{ display: flex; gap: 0.75rem; flex-wrap: wrap; }}
        .field {{ flex: 1 1 200px; display: flex; flex-direction: column; gap: 0.2rem; }}
        label {{ font-size: 0.8rem; font-weight: 600; color: #4b5563; }}
        textarea, select {{ width: 100%; padding: 0.45rem 0.55rem; border-radius: 10px; border: 1px solid #d1d5db; font-size: 0.9rem; font-family: inherit; }}
        textarea {{ resize: vertical; min-height: 70px; }}
        textarea:focus, select:focus {{ outline: none; border-color: #6366f1; box-shadow: 0 0 0 1px #6366f1; }}
        .actions {{ display: flex; align-items: center; justify-content: space-between; gap: 0.75rem; margin-top: 0.25rem; }}
        button[type="submit"] {{ padding: 0.6rem 1.2rem; border-radius: 999px; border: none; background: #111827; color: white; font-size: 0.9rem; font-weight: 500; cursor: pointer; display: inline-flex; align-items: center; gap: 0.4rem; }}
        button[type="submit"]:disabled {{ opacity: 0.6; cursor: default; }}
        .btn-ppt {{ padding: 0.6rem 1.2rem; border-radius: 999px; border: none; background: #4f46e5; color: white; font-size: 0.9rem; font-weight: 500; cursor: pointer; display: inline-flex; align-items: center; gap: 0.4rem; }}
        .btn-ppt:disabled {{ opacity: 0.6; cursor: default; }}
        .muted {{ color: #6b7280; font-size: 0.75rem; }}
        .error-text {{ color: #b91c1c; font-size: 0.8rem; margin-top: 0.25rem; }}
        @media (max-width: 640px) {{ .app {{ padding: 0.75rem; }} .chat-window {{ max-height: 55vh; }} .bubble {{ max-width: 90%; }} }}
      </style>
    </head>
    <body>
      <div class="app">
        <div class="card">
          <header class="header">
            <div class="title-row">
              <div class="pill">L</div>
              <div style="flex:1;">
                <h1>LIFT: Lecture and Instructional Faculty Tool
 </h1>
                <p class="subtitle">Chat with LIFT using course content + a prompt, or generate a PowerPoint from your prompt.</p>
              </div>
              <button id="reset-btn" title="Start a new conversation" style="background:none;border:1px solid #e5e7eb;border-radius:999px;padding:0.3rem 0.75rem;font-size:0.78rem;color:#6b7280;cursor:pointer;white-space:nowrap;">&#8635; New Chat</button>
            </div>
          </header>

          <main class="body">
            <div id="chat-window" class="chat-window">
              <div class="message assistant">
                <div class="avatar">L</div>
                <div class="bubble">
                  <div class="name">LIFT</div>
                  <div class="bubble-body">Hi! Select a use case, enter a prompt, and optionally upload a .txt or .docx file. Click <strong>Generate with LIFT</strong> for a text response, or <strong>Generate as PowerPoint</strong> to download a .pptx.</div>
                </div>
              </div>
            </div>

            <!-- ── Chat / Text Generation Form ── -->
            <div class="section-title">Generate Teaching Materials</div>
            <form id="lift-form" action="/generate-content" method="POST" enctype="multipart/form-data" class="form-wrapper">
              <div class="form-row">
                <div class="field">
                  <label for="use_case">Select a Use Case</label>
                  <select id="use_case" name="use_case">
                    <option value="none">General / No Specific Use Case</option>
                    <option value="uc1">Use Case 1: Rapid Course Material Development</option>
                    <option value="uc2">Use Case 2: Accessible Multi-Format Materials</option>
                    <option value="uc3">Use Case 3: Assessment &amp; Rubric Development</option>
                    <option value="uc4">Use Case 4: Flipped Classroom Content</option>
                    <option value="uc5">Use Case 5: Cross-Disciplinary Revision</option>
                    <option value="uc6">Use Case 6: Lecture Presentation Creator</option>
                    <option value="uc7">Use Case 7: Faculty AI Literacy &amp; Development</option>
                    <option value="uc8">Use Case 8: Cross-Campus Collaborative Design</option>
                    <option value="uc9">Use Case 9: Pedagogical Research Design</option>
                  </select>
                </div>
                <div class="field">
                  <label for="file">Upload content file (.txt or .docx)</label>
                  <input id="file" type="file" name="file" accept=".txt,.docx" />
                </div>
              </div>
              <div class="form-row">
                <div class="field">
                  <label for="instructions">Prompt</label>
                  <textarea id="instructions" name="instructions" placeholder="e.g., Write 5 quiz questions..."></textarea>
                </div>
              </div>
              <div class="actions">
                <div class="muted">Conversation memory is enabled per browser.</div>
                <button type="submit"><span>Generate with LIFT</span></button>
              </div>
              <div id="error" class="error-text" style="display:none;"></div>
              <div id="ppt-status" class="muted" style="display:none;margin-top:0.25rem;"></div>
              <div style="text-align:right;margin-top:0.4rem;display:flex;flex-direction:column;align-items:flex-end;gap:0.25rem;">
                <button type="button" id="docx-btn" style="background:none;border:none;padding:0;color:#9ca3af;font-size:0.75rem;cursor:pointer;text-decoration:underline;">generate as docx</button>
                <button type="button" id="ppt-btn" style="background:none;border:none;padding:0;color:#9ca3af;font-size:0.75rem;cursor:pointer;text-decoration:underline;">generate as powerpoint</button>
              </div>
            </form>
          </main>

        </div>
      </div>

      <script>
        const form = document.getElementById('lift-form');
        const chat = document.getElementById('chat-window');
        const errorBox = document.getElementById('error');
        const submitBtn = form.querySelector('button[type="submit"]');

        function escapeHTML(str) {{ return str.replace(/&/g, "&amp;").replace(/</g, "&lt;").replace(/>/g, "&gt;").replace(/"/g, "&quot;").replace(/'/g, "&#039;"); }}
        function formatWithBreaks(str) {{
          let s = escapeHTML(str);
          s = s.replace(/\*\*(.*?)\*\*/gs, '<strong>$1</strong>');
          s = s.replace(/\*(.*?)\*/gs, '$1');
          s = s.replace(/^#+\s?/gm, '');
          s = s.replace(/\\n/g, "<br>");
          return s;
        }}
        function scrollChatToBottom() {{ chat.scrollTop = chat.scrollHeight; }}

        function addMessage(role, htmlBody) {{
          const wrapper = document.createElement('div');
          wrapper.className = 'message ' + role;
          const avatar = document.createElement('div');
          avatar.className = 'avatar';
          avatar.textContent = role === 'user' ? 'You' : 'L';
          const bubble = document.createElement('div');
          bubble.className = 'bubble';
          const name = document.createElement('div');
          name.className = 'name';
          name.textContent = role === 'user' ? 'You' : 'LIFT';
          const body = document.createElement('div');
          body.className = 'bubble-body';
          body.innerHTML = htmlBody;
          bubble.appendChild(name);
          bubble.appendChild(body);
          wrapper.appendChild(avatar);
          wrapper.appendChild(bubble);
          chat.appendChild(wrapper);
          scrollChatToBottom();
          return wrapper;
        }}

        // ── Chat form handler ──
        form.addEventListener('submit', async (e) => {{
          e.preventDefault();
          errorBox.style.display = 'none';
          const instructions = document.getElementById('instructions').value.trim();
          const useCaseSelect = document.getElementById('use_case');
          const useCaseText = useCaseSelect.options[useCaseSelect.selectedIndex].text;
          const fileInput = document.getElementById('file');
          const file = fileInput.files[0];

          if (!instructions && !file) {{
            errorBox.textContent = 'Provide a prompt or upload a file.';
            errorBox.style.display = 'block';
            return;
          }}

          let summaryParts = [`<strong>Mode:</strong> ${{useCaseText}}`];
          if (instructions) summaryParts.push('<strong>Prompt:</strong><br>' + formatWithBreaks(instructions));
          if (file) summaryParts.push('<strong>File:</strong> ' + escapeHTML(file.name));

          addMessage('user', summaryParts.join('<br><br>'));
          const typingMsg = addMessage('assistant', '<span class="typing-dots"><span class="dot"></span><span class="dot"></span><span class="dot"></span></span>');
          submitBtn.disabled = true;

          try {{
            const res = await fetch('/generate-content', {{ method: 'POST', body: new FormData(form) }});
            const data = await res.json();
            const bubbleBody = typingMsg.querySelector('.bubble-body');
            if (!res.ok) bubbleBody.innerHTML = '<span class="error-text">Error: ' + (data.error || 'Failed') + '</span>';
            else bubbleBody.innerHTML = formatWithBreaks(data.generated_text);
          }} catch (err) {{
            typingMsg.querySelector('.bubble-body').innerHTML = 'Network error.';
          }} finally {{
            submitBtn.disabled = false;
            scrollChatToBottom();
          }}
        }});

        // ── Docx button handler ──
        document.getElementById('docx-btn').addEventListener('click', async () => {{
          const pptStatus = document.getElementById('ppt-status');
          const docxBtn = document.getElementById('docx-btn');
          const instructions = document.getElementById('instructions').value.trim();
          const file = document.getElementById('file').files[0];
          errorBox.style.display = 'none';

          if (!instructions && !file) {{
            errorBox.textContent = 'Provide a prompt or upload a file before generating a Word document.';
            errorBox.style.display = 'block';
            return;
          }}

          docxBtn.disabled = true;
          submitBtn.disabled = true;
          pptStatus.textContent = 'Generating your Word document\u2026 this may take up to 60 seconds.';
          pptStatus.style.display = 'block';

          try {{
            const res = await fetch('/generate-docx', {{ method: 'POST', body: new FormData(form) }});
            if (!res.ok) {{
              const data = await res.json();
              errorBox.textContent = 'Docx Error: ' + (data.error || 'Failed to generate.');
              errorBox.style.display = 'block';
              pptStatus.style.display = 'none';
            }} else {{
              const blob = await res.blob();
              const disposition = res.headers.get('Content-Disposition') || '';
              const match = disposition.match(/filename="(.+?)"/);
              const filename = match ? match[1] : 'LIFT_Document.docx';
              const url = URL.createObjectURL(blob);
              const a = document.createElement('a');
              a.href = url; a.download = filename;
              document.body.appendChild(a); a.click(); a.remove();
              URL.revokeObjectURL(url);
              pptStatus.textContent = '\u2705 Word document downloaded!';
            }}
          }} catch (err) {{
            errorBox.textContent = 'Network error. Please try again.';
            errorBox.style.display = 'block';
            pptStatus.style.display = 'none';
          }} finally {{
            docxBtn.disabled = false;
            submitBtn.disabled = false;
          }}
        }});

        // ── PPT button handler (uses the same main form fields) ──
        document.getElementById('ppt-btn').addEventListener('click', async () => {{
          const pptStatus = document.getElementById('ppt-status');
          const pptBtn = document.getElementById('ppt-btn');
          const instructions = document.getElementById('instructions').value.trim();
          const fileInput = document.getElementById('file');
          const file = fileInput.files[0];
          errorBox.style.display = 'none';

          if (!instructions && !file) {{
            errorBox.textContent = 'Provide a prompt or upload a file before generating a PowerPoint.';
            errorBox.style.display = 'block';
            return;
          }}

          pptBtn.disabled = true;
          submitBtn.disabled = true;
          pptStatus.textContent = 'Generating your PowerPoint\u2026 this may take up to 60 seconds.';
          pptStatus.style.display = 'block';

          try {{
            const res = await fetch('/generate-ppt', {{ method: 'POST', body: new FormData(form) }});
            if (!res.ok) {{
              const data = await res.json();
              errorBox.textContent = 'PPT Error: ' + (data.error || 'Failed to generate.');
              errorBox.style.display = 'block';
              pptStatus.style.display = 'none';
            }} else {{
              const blob = await res.blob();
              const disposition = res.headers.get('Content-Disposition') || '';
              const match = disposition.match(/filename="(.+?)"/);
              const filename = match ? match[1] : 'LIFT_Presentation.pptx';
              const url = URL.createObjectURL(blob);
              const a = document.createElement('a');
              a.href = url; a.download = filename;
              document.body.appendChild(a); a.click(); a.remove();
              URL.revokeObjectURL(url);
              pptStatus.textContent = '\u2705 PowerPoint downloaded!';
            }}
          }} catch (err) {{
            errorBox.textContent = 'Network error. Please try again.';
            errorBox.style.display = 'block';
            pptStatus.style.display = 'none';
          }} finally {{
            pptBtn.disabled = false;
            submitBtn.disabled = false;
          }}
        }});
        // ── Reset / New Chat button ──
        document.getElementById('reset-btn').addEventListener('click', async () => {{
          await fetch('/reset', {{ method: 'POST' }});
          chat.innerHTML = '';
          addMessage('assistant', 'Conversation cleared. Select a use case, enter a prompt, and let\u2019s start fresh!');
          document.getElementById('instructions').value = '';
          document.getElementById('file').value = '';
          errorBox.style.display = 'none';
          document.getElementById('ppt-status').style.display = 'none';
        }});
      </script>
    </body>
    </html>
    """
    return Response(html, mimetype="text/html")

@app.route("/healthz", methods=["GET"])
def healthz():
    return jsonify({"status": "ok"})

@app.route("/reset", methods=["POST"])
def reset():
    session.clear()
    return jsonify({"status": "cleared"})

def _get_history():
    return session.get("lift_history", [])

def _save_history(history):
    session["lift_history"] = history[-MAX_HISTORY_TURNS:]

def _build_history_block(history):
    if not history: return "No prior conversation.\n"
    return "\n".join([f"{turn['role'].upper()}:\n{turn['content']}\n" for turn in history])

def _extract_file_text(uploaded_file) -> str:
    """Extract plain text from .txt or .docx uploads."""
    filename = uploaded_file.filename.lower()
    raw = uploaded_file.read()
    if filename.endswith(".docx"):
        doc = Document(BytesIO(raw))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    else:
        return raw.decode("utf-8", errors="ignore")

@app.route("/generate-content", methods=["POST"])
def generate_content():
    instructions = request.form.get("instructions", "") or ""
    use_case_key = request.form.get("use_case", "none")
    uploaded_file = request.files.get("file")

    use_case_context = USE_CASES.get(use_case_key, USE_CASES["none"])

    combined_text = ""
    if uploaded_file and uploaded_file.filename:
        try:
            combined_text += _extract_file_text(uploaded_file) + "\n"
        except Exception as e:
            return jsonify({"error": f"File error: {e}"}), 400

    history = _get_history()
    history_block = _build_history_block(history)

    # Scrub PII from user-supplied text before sending to the LLM
    safe_instructions = scrub_pii(instructions)
    safe_content = scrub_pii(combined_text)

    prompt = f"""You are LIFT, an AI assistant for faculty.

STRICT OPERATING CONTEXT:
{use_case_context}

General Capabilities:
- Learning outcomes & scaffolding
- Summary & Quiz generation (Bloom's alignment)
- Accessibility & Flipped classroom materials

===== PRIOR CONVERSATION =====
{history_block}
===== END PRIOR CONVERSATION =====

LATEST REQUEST:
Prompt: {safe_instructions}
Content: {safe_content}

Respond as LIFT using the specific Use Case context provided above.
"""

    try:
        resp = model.generate_content(prompt, request_options={"timeout": 110})
        output_text = getattr(resp, "text", "")

        history.append({"role": "user", "content": f"Use Case: {use_case_key} | Prompt: {instructions[:200]}"})
        history.append({"role": "assistant", "content": output_text[:ASSISTANT_SNIPPET_CHARS]})
        _save_history(history)

        return jsonify({"generated_text": output_text})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/generate-docx", methods=["POST"])
def generate_docx():
    """
    Same form fields as /generate-content. Returns Gemini's output as a formatted .docx.
    """
    instructions = request.form.get("instructions", "") or ""
    use_case_key = request.form.get("use_case", "none")
    uploaded_file = request.files.get("file")

    if not instructions and (not uploaded_file or not uploaded_file.filename):
        return jsonify({"error": "Provide a prompt or upload a file to generate a Word document."}), 400

    use_case_context = USE_CASES.get(use_case_key, USE_CASES["none"])

    combined_text = ""
    if uploaded_file and uploaded_file.filename:
        try:
            combined_text = _extract_file_text(uploaded_file) + "\n"
        except Exception as e:
            return jsonify({"error": f"File error: {e}"}), 400

    safe_instructions = scrub_pii(instructions)
    safe_content = scrub_pii(combined_text)

    prompt = f"""You are LIFT, an AI assistant for faculty.

STRICT OPERATING CONTEXT:
{use_case_context}

General Capabilities:
- Learning outcomes & scaffolding
- Summary & Quiz generation (Bloom's alignment)
- Accessibility & Flipped classroom materials

LATEST REQUEST:
Prompt: {safe_instructions}
Content: {safe_content}

Respond as LIFT using the specific Use Case context provided above.
"""

    try:
        resp = model.generate_content(prompt, request_options={"timeout": 110})
        output_text = getattr(resp, "text", "")
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    # ── Build the Word document ───────────────────────────────
    def _add_bold_runs(para, text):
        """Split text on **bold** markers and add runs with correct bold state."""
        # Remove any lone/unmatched single asterisks first
        text = re.sub(r'(?<!\*)\*(?!\*)', '', text)
        parts = re.split(r'\*\*(.*?)\*\*', text)
        for i, part in enumerate(parts):
            if part:
                run = para.add_run(part)
                run.bold = (i % 2 == 1)

    try:
        doc = Document()

        # Title: first non-empty line with all asterisks/hashes stripped
        first_line = output_text.strip().splitlines()[0] if output_text.strip() else "LIFT Output"
        title_text = re.sub(r'[#*]', '', first_line).strip()[:120]
        doc.add_heading(title_text, level=0)

        for line in output_text.splitlines():
            stripped = line.strip()
            if not stripped:
                doc.add_paragraph("")
                continue

            # ── Markdown headings → Word headings ──
            if stripped.startswith("### "):
                doc.add_heading(re.sub(r'\*', '', stripped[4:]).strip(), level=3)
            elif stripped.startswith("## "):
                doc.add_heading(re.sub(r'\*', '', stripped[3:]).strip(), level=2)
            elif stripped.startswith("# "):
                doc.add_heading(re.sub(r'\*', '', stripped[2:]).strip(), level=1)

            # ── Entire line is **bold** → treat as a sub-heading ──
            elif re.match(r'^\*\*.+\*\*$', stripped):
                doc.add_heading(re.sub(r'\*', '', stripped).strip(), level=2)

            # ── Bullet points (with inline bold support) ──
            elif stripped.startswith("* ") or stripped.startswith("- "):
                para = doc.add_paragraph(style="List Bullet")
                _add_bold_runs(para, stripped[2:])

            # ── Regular paragraph with inline bold ──
            else:
                para = doc.add_paragraph()
                _add_bold_runs(para, stripped)

        output = BytesIO()
        doc.save(output)
        output.seek(0)

        safe_title = re.sub(r'[^\w\s\-]', '', title_text)
        filename = safe_title.strip().replace(" ", "_")[:50] + ".docx"

        return Response(
            output.getvalue(),
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )

    except Exception as e:
        return jsonify({"error": f"Docx build error: {e}"}), 500


@app.route("/generate-ppt", methods=["POST"])
def generate_ppt():
    """
    Uses the same form fields as /generate-content (use_case, instructions, file).
    Asks Gemini to produce structured slide JSON, then builds and returns a .pptx.
    """
    instructions = request.form.get("instructions", "") or ""
    use_case_key = request.form.get("use_case", "none")
    uploaded_file = request.files.get("file")

    if not instructions and (not uploaded_file or not uploaded_file.filename):
        return jsonify({"error": "Provide a prompt or upload a file to generate a PowerPoint."}), 400

    use_case_context = USE_CASES.get(use_case_key, USE_CASES["none"])

    combined_text = ""
    if uploaded_file and uploaded_file.filename:
        try:
            combined_text = _extract_file_text(uploaded_file) + "\n"
        except Exception as e:
            return jsonify({"error": f"File error: {e}"}), 400

    safe_instructions = scrub_pii(instructions)
    safe_content = scrub_pii(combined_text)

    ppt_prompt = f"""You are LIFT, an AI assistant for faculty. Generate a complete, pedagogically-structured PowerPoint presentation based on the request below.

PEDAGOGICAL CONTEXT:
{use_case_context}

FACULTY REQUEST / PROMPT:
{safe_instructions if safe_instructions else "No specific prompt provided — generate a comprehensive presentation from the content below."}

REFERENCE CONTENT (if provided):
{safe_content[:12000] if safe_content.strip() else "None."}

Return ONLY valid JSON — no markdown fences, no explanation, no extra text. Use this exact schema:
{{
  "title": "Presentation Title",
  "subtitle": "Optional subtitle or course name",
  "slides": [
    {{
      "title": "Slide Title",
      "bullets": ["Key point 1", "Key point 2", "Key point 3"],
      "notes": "Full presenter notes with learning-theory rationale and timing guidance."
    }}
  ]
}}

Guidelines:
- Create 10–15 slides that directly address the faculty's prompt.
- Max 5 bullets per slide; each bullet is a concise phrase, not a full sentence.
- Apply Gagné's Nine Events: attention → objectives → recall → content → practice → feedback → assess → transfer.
- First slide = title/introduction; last slide = summary + discussion questions.
- Presenter notes: full narration script, active-learning pause points, common student questions.
"""

    try:
        resp = model.generate_content(ppt_prompt, request_options={"timeout": 110})
        raw = (getattr(resp, "text", "") or "").strip()

        # Strip any accidental markdown code fences
        if raw.startswith("```"):
            parts = raw.split("```")
            raw = parts[1] if len(parts) > 1 else raw
            if raw.startswith("json"):
                raw = raw[4:]
            raw = raw.strip()

        slide_data = json.loads(raw)
    except json.JSONDecodeError as e:
        return jsonify({"error": f"AI returned malformed JSON: {e}. Try again or simplify your document."}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    # ── Build the PowerPoint ──────────────────────────────────
    try:
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        INDIGO = RGBColor(0x4F, 0x46, 0xE5)
        WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
        DARK   = RGBColor(0x11, 0x18, 0x27)

        def _set_font(run, size_pt, bold=False, color=None):
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color

        # Title slide (layout 0)
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.background.fill.solid()
        title_slide.background.fill.fore_color.rgb = INDIGO

        ts_title = title_slide.shapes.title
        ts_title.text = slide_data.get("title", "LIFT Presentation")
        for para in ts_title.text_frame.paragraphs:
            for run in para.runs:
                _set_font(run, 36, bold=True, color=WHITE)

        if len(title_slide.placeholders) > 1:
            ts_sub = title_slide.placeholders[1]
            ts_sub.text = slide_data.get("subtitle", "")
            for para in ts_sub.text_frame.paragraphs:
                for run in para.runs:
                    _set_font(run, 20, color=RGBColor(0xC7, 0xD2, 0xFE))

        # Content slides (layout 1 — title + content)
        for slide_info in slide_data.get("slides", []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])

            # Title
            sh_title = slide.shapes.title
            sh_title.text = slide_info.get("title", "")
            for para in sh_title.text_frame.paragraphs:
                for run in para.runs:
                    _set_font(run, 24, bold=True, color=INDIGO)

            # Bullets
            bullets = slide_info.get("bullets", [])
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        tf.paragraphs[0].text = bullet
                        para = tf.paragraphs[0]
                    else:
                        para = tf.add_paragraph()
                        para.text = bullet
                    para.level = 0
                    for run in para.runs:
                        _set_font(run, 18, color=DARK)

            # Presenter notes
            notes_text = slide_info.get("notes", "")
            if notes_text:
                slide.notes_slide.notes_text_frame.text = notes_text

        # Save
        output = BytesIO()
        prs.save(output)
        output.seek(0)

        safe_title = re.sub(r'[^\w\s\-]', '', slide_data.get("title", "LIFT_Presentation"))
        filename = safe_title.strip().replace(" ", "_")[:50] + ".pptx"

        return Response(
            output.getvalue(),
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )

    except Exception as e:
        return jsonify({"error": f"PPT build error: {e}"}), 500


if __name__ == "__main__":
    port = int(os.getenv("PORT", "8080"))
    app.run(host="0.0.0.0", port=port)
